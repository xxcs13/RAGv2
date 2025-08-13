"""
Document parsing system for PDF, PPTX, and Excel files.
"""
import re
import os
import pandas as pd
from typing import List, Dict, Any
from pathlib import Path
from pptx import Presentation
from pypdf import PdfReader
import pdfplumber


class PDFParser:
    """PDF text extraction with number formatting correction."""
    
    def parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract text and tables from PDF file."""
        try:
            filename = Path(file_path).name
            pages = []
            
            try:
                pages = self._extract_with_pdfplumber(file_path)
                print(f"Successfully parsed PDF with pdfplumber: {filename}")
            except Exception as e:
                print(f"pdfplumber failed, trying pypdf: {e}")
                try:
                    pages = self._extract_with_pypdf(file_path)
                    print(f"Successfully parsed PDF with pypdf: {filename}")
                except Exception as e2:
                    print(f"Both PDF extraction methods failed: {e2}")
                    return self._create_fallback_report(file_path)
            
            processed_pages = []
            for page_data in pages:
                processed_text = self._post_process_text(page_data['text'])
                if processed_text.strip():
                    processed_pages.append({
                        'page': page_data['page'],
                        'text': processed_text.strip()
                    })
            
            report = {
                'metainfo': {
                    'sha1_name': filename.rsplit('.', 1)[0],
                    'filename': filename,
                    'pages_amount': len(processed_pages),
                    'text_blocks_amount': len(processed_pages),
                    'tables_amount': 0,
                    'pictures_amount': 0,
                    'document_type': 'pdf'
                },
                'content': {'pages': processed_pages},
                'tables': [],
                'pictures': []
            }
            
            print(f"Successfully parsed PDF: {filename} ({len(processed_pages)} pages)")
            return report
            
        except Exception as e:
            print(f"Error parsing PDF file {file_path}: {e}")
            return self._create_fallback_report(file_path)
    
    def _extract_with_pdfplumber(self, file_path: str) -> List[Dict]:
        """Extract text and tables using pdfplumber with layout detection."""
        pages = []
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text_parts = []
                
                # Perform layout analysis on the page
                layout_info = self._analyze_page_layout(page)
                
                # Extract text based on layout type
                if layout_info['layout_type'] == 'multi_column':
                    # Multi-column extraction with coordinate-based column detection
                    column_texts = self._extract_multi_column_text(page, layout_info)
                    if column_texts:
                        text_parts.extend(column_texts)
                else:
                    # Standard single-column extraction
                    standard_text = page.extract_text()
                    if standard_text:
                        text_parts.append(standard_text)
                
                # Extract layout-preserved text for complex layouts
                try:
                    layout_text = page.extract_text(layout=True, x_tolerance=1, y_tolerance=1)
                    if layout_text and layout_text not in text_parts:
                        text_parts.append("=== Layout Preserved ===")
                        text_parts.append(layout_text)
                except:
                    pass
                
                # Extract tables
                try:
                    tables = page.extract_tables()
                    if tables:
                        text_parts.append("=== Tables ===")
                        for i, table in enumerate(tables):
                            if table:
                                table_text = self._format_table_text(table)
                                text_parts.append(f"Table {i+1}:\n{table_text}")
                except:
                    pass
                
                combined_text = '\n\n'.join(text_parts)
                if combined_text.strip():
                    page_data = {
                        'page': page_num,
                        'text': combined_text.strip(),
                        'layout_type': layout_info['layout_type'],
                        'column_count': layout_info['column_count'],
                        'extraction_method': layout_info['extraction_method']
                    }
                    pages.append(page_data)
        
        return pages
    
    def _analyze_page_layout(self, page) -> Dict[str, Any]:
        """Analyze PDF page layout to detect columns and structure."""
        try:
            # Get page dimensions
            page_width = page.width
            page_height = page.height
            
            # Extract characters with coordinates for layout analysis
            chars = page.chars
            if not chars:
                return {
                    'layout_type': 'single_column',
                    'column_count': 1,
                    'extraction_method': 'standard'
                }
            
            # Analyze character distribution across page width
            x_positions = [char['x0'] for char in chars if 'x0' in char]
            if not x_positions:
                return {
                    'layout_type': 'single_column',
                    'column_count': 1,
                    'extraction_method': 'standard'
                }
            
            # Group characters by approximate x-position to detect columns
            x_bins = self._create_position_bins(x_positions, page_width)
            column_gaps = self._detect_column_gaps(x_bins, page_width)
            
            # Determine layout type based on gaps and distribution
            if len(column_gaps) >= 1:
                # Check if gaps are significant enough to indicate columns
                significant_gaps = [gap for gap in column_gaps if gap['width'] > page_width * 0.02]
                
                if significant_gaps:
                    column_count = len(significant_gaps) + 1
                    layout_type = 'multi_column' if column_count > 1 else 'single_column'
                else:
                    column_count = 1
                    layout_type = 'single_column'
            else:
                column_count = 1
                layout_type = 'single_column'
            
            # Additional complexity check based on text object distribution
            text_objects = page.objects.get('char', [])
            if len(text_objects) > 1000 and layout_type == 'single_column':
                layout_type = 'complex'
            
            return {
                'layout_type': layout_type,
                'column_count': column_count,
                'extraction_method': 'coordinate_based',
                'page_width': page_width,
                'column_gaps': column_gaps if 'significant_gaps' in locals() else []
            }
            
        except Exception as e:
            print(f"Warning: Layout analysis failed for page, using default: {e}")
            return {
                'layout_type': 'single_column',
                'column_count': 1,
                'extraction_method': 'standard'
            }
    
    def _create_position_bins(self, x_positions: List[float], page_width: float, bin_count: int = 50) -> List[int]:
        """Create bins for x-position distribution analysis."""
        if not x_positions:
            return []
        
        # Create histogram bins across page width
        bin_width = page_width / bin_count
        bins = [0] * bin_count
        
        for x in x_positions:
            bin_index = min(int(x / bin_width), bin_count - 1)
            bins[bin_index] += 1
        
        return bins
    
    def _detect_column_gaps(self, x_bins: List[int], page_width: float) -> List[Dict]:
        """Detect gaps between columns based on character distribution."""
        if not x_bins:
            return []
        
        # Find regions with very low character density (potential gaps)
        total_chars = sum(x_bins)
        if total_chars == 0:
            return []
        
        # Calculate moving average to smooth the distribution
        window_size = 3
        smoothed_bins = []
        for i in range(len(x_bins)):
            start_idx = max(0, i - window_size // 2)
            end_idx = min(len(x_bins), i + window_size // 2 + 1)
            avg = sum(x_bins[start_idx:end_idx]) / (end_idx - start_idx)
            smoothed_bins.append(avg)
        
        # Identify gaps (regions with density below threshold)
        avg_density = total_chars / len(x_bins)
        gap_threshold = avg_density * 0.1  # 10% of average density
        
        gaps = []
        in_gap = False
        gap_start = 0
        bin_width = page_width / len(x_bins)
        
        for i, density in enumerate(smoothed_bins):
            if density <= gap_threshold and not in_gap:
                # Start of gap
                in_gap = True
                gap_start = i
            elif density > gap_threshold and in_gap:
                # End of gap
                in_gap = False
                gap_width = (i - gap_start) * bin_width
                if gap_width > page_width * 0.02:  # Minimum 2% of page width
                    gaps.append({
                        'start': gap_start * bin_width,
                        'end': i * bin_width,
                        'width': gap_width
                    })
        
        # Handle gap that extends to end of page
        if in_gap:
            gap_width = (len(smoothed_bins) - gap_start) * bin_width
            if gap_width > page_width * 0.02:
                gaps.append({
                    'start': gap_start * bin_width,
                    'end': page_width,
                    'width': gap_width
                })
        
        return gaps
    
    def _extract_multi_column_text(self, page, layout_info: Dict) -> List[str]:
        """Extract text from multi-column layout using coordinate-based approach."""
        try:
            column_gaps = layout_info.get('column_gaps', [])
            page_width = layout_info.get('page_width', page.width)
            
            if not column_gaps:
                # Fallback to standard extraction
                standard_text = page.extract_text()
                return [standard_text] if standard_text else []
            
            # Define column boundaries based on detected gaps
            column_boundaries = [0]  # Start with left edge
            for gap in column_gaps:
                column_boundaries.append(gap['start'])
                column_boundaries.append(gap['end'])
            column_boundaries.append(page_width)  # End with right edge
            
            # Remove duplicates and sort
            column_boundaries = sorted(list(set(column_boundaries)))
            
            # Create column regions (pairs of boundaries)
            columns = []
            for i in range(0, len(column_boundaries) - 1, 2):
                if i + 1 < len(column_boundaries):
                    left = column_boundaries[i]
                    right = column_boundaries[i + 1]
                    columns.append((left, right))
            
            # Extract text from each column
            column_texts = []
            for i, (left, right) in enumerate(columns):
                # Create a cropped version of the page for this column
                cropped_page = page.within_bbox((left, 0, right, page.height))
                column_text = cropped_page.extract_text()
                
                if column_text and column_text.strip():
                    column_texts.append(f"||COLUMN|| {column_text}")
            
            return column_texts if column_texts else [page.extract_text()]
            
        except Exception as e:
            print(f"Warning: Multi-column extraction failed, using standard: {e}")
            standard_text = page.extract_text()
            return [standard_text] if standard_text else []
    
    def _extract_with_pypdf(self, file_path: str) -> List[Dict]:
        """Extract text using pypdf as fallback."""
        pages = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    pages.append({'page': page_num, 'text': text.strip()})
        
        return pages
    
    def _format_table_text(self, table: List[List]) -> str:
        """Format extracted table data into readable text."""
        if not table:
            return ""
        
        formatted_rows = []
        for row in table:
            if row:
                clean_cells = []
                for cell in row:
                    cell_text = str(cell).strip() if cell is not None else ""
                    clean_cells.append(cell_text)
                formatted_rows.append(" | ".join(clean_cells))
        
        return "\n".join(formatted_rows)
    
    def _post_process_text(self, text: str) -> str:
        """Fix number formatting and spacing issues in extracted text."""
        if not text:
            return text
        
        processed_text = self._fix_number_formatting(text)
        processed_text = self._fix_spacing_issues(processed_text)
        processed_text = self._fix_symbols(processed_text)
        
        return processed_text
    
    def _fix_number_formatting(self, text: str) -> str:
        """Correct decimal points, thousand separators, and currency symbols."""
        # Fix thousand separators (including Chinese comma)
        text = re.sub(r'(\d)\s*[,\uff0c]\s*(\d{3})', r'\1,\2', text)
        # Fix decimal points (including Chinese period)
        text = re.sub(r'(\d)\s*[.\uff0e]\s*(\d)', r'\1.\2', text)
        # Fix number spacing
        text = re.sub(r'(\d)\s*[oO]\s*(\d)', r'\1.0\2', text)
        text = re.sub(r'(\d)\s+(\d)', r'\1\2', text)
        # Fix percentage symbols
        text = re.sub(r'(\d)\s*[%\uff05]', r'\1%', text)
        # Fix currency symbols
        text = re.sub(r'[$\uff04]\s*(\d)', r'$\1', text)
        text = re.sub(r'([NT$]+)\s*(\d)', r'\1\2', text)
        # Fix negative signs
        text = re.sub(r'[-\uff0d\u2014]\s*(\d)', r'-\1', text)
        
        return text
    
    def _fix_spacing_issues(self, text: str) -> str:
        """Fix spacing and line break issues."""
        text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text)
        text = re.sub(r' {2,}', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        return text
    
    def _fix_symbols(self, text: str) -> str:
        """Fix currency and other symbols."""
        # Normalize currency symbols
        text = re.sub(r'[$\uff04]', '$', text)
        # Normalize percentage symbols
        text = re.sub(r'[%\uff05]', '%', text)
        # Normalize parentheses
        text = re.sub(r'[\uff08(]', '(', text)
        text = re.sub(r'[\uff09)]', ')', text)
        
        return text
    
    def _create_fallback_report(self, file_path: str) -> Dict[str, Any]:
        """Create minimal report when PDF parsing fails."""
        filename = Path(file_path).name
        return {
            'metainfo': {
                'sha1_name': filename.rsplit('.', 1)[0],
                'filename': filename,
                'pages_amount': 0,
                'text_blocks_amount': 0,
                'tables_amount': 0,
                'pictures_amount': 0,
                'document_type': 'failed'
            },
            'content': {'pages': []},
            'tables': [],
            'pictures': []
        }


class PPTXParser:
    """Enhanced PPTX content extraction with structured data conversion.
    
    This parser extracts:
    1. Text content as usual
    2. Tables converted to CSV format for ExcelParser processing
    3. Charts converted to CSV format for ExcelParser processing
    
    The structured data (tables/charts) are processed by AdaptiveExcelParser
    to maintain data integrity and enable better LLM understanding.
    """
    
    def __init__(self):
        """Initialize with ExcelParser for structured data processing."""
        self.excel_parser = ExcelParser()
    
    def parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """Parse PPTX file with robust structured data extraction and CSV conversion."""
        try:
            filename = Path(file_path).name
            
            # Extract structured content using robust methodology
            structured_data = self._extract_structured_content_robust(file_path)
            
            # Process the extracted content
            pages = []
            all_csv_data = []
            tables_found = 0
            charts_found = 0
            
            for slide_data in structured_data['slides']:
                slide_num = slide_data['index']
                slide_content = []
                
                # Add slide title
                if slide_data.get('title'):
                    slide_content.append(f"=== SLIDE {slide_num} TITLE ===")
                    slide_content.append(slide_data['title'])
                    slide_content.append("")
                
                # Add text blocks (unstructured content)
                text_content = []
                for text_block in slide_data.get('text_blocks', []):
                    if text_block.strip():
                        text_content.append(text_block.strip())
                
                if text_content:
                    slide_content.append("=== TEXT CONTENT ===")
                    slide_content.extend(text_content)
                    slide_content.append("")
                
                # Process tables as CSV and combine with text
                for i, table_csv in enumerate(slide_data.get('table_csvs', []), 1):
                    tables_found += 1
                    if table_csv.strip():
                        # Process CSV with ExcelParser
                        processed_table = self._process_csv_with_excel_parser(
                            table_csv, f"Slide {slide_num} Table {i}"
                        )
                        slide_content.append(f"=== TABLE {i} DATA ===")
                        slide_content.append(processed_table)
                        slide_content.append("")
                
                # Process charts as CSV and combine with text  
                for i, chart_csv in enumerate(slide_data.get('chart_csvs', []), 1):
                    charts_found += 1
                    if chart_csv.strip():
                        # Process CSV with ExcelParser
                        processed_chart = self._process_csv_with_excel_parser(
                            chart_csv, f"Slide {slide_num} Chart {i}"
                        )
                        slide_content.append(f"=== CHART {i} DATA ===")
                        slide_content.append(processed_chart)
                        slide_content.append("")
                
                # Combine all content for this slide
                if slide_content:
                    pages.append({
                        'page': slide_num,
                        'text': '\n'.join(slide_content)
                    })
            
            # Create report structure
            report = {
                'metainfo': {
                    'sha1_name': filename.rsplit('.', 1)[0],
                    'filename': filename,
                    'pages_amount': len(pages),
                    'text_blocks_amount': len(pages),
                    'tables_amount': tables_found,
                    'pictures_amount': 0,
                    'document_type': 'pptx'
                },
                'content': {'pages': pages},
                'tables': [],
                'pictures': []
            }
            
            print(f"Successfully parsed PPTX: {filename}")
            print(f"  - {len(pages)} slides with content")
            print(f"  - {tables_found} tables, {charts_found} charts processed as structured data")
            print(f"  - Total structured data items: {tables_found + charts_found}")
            
            return report
            
        except Exception as e:
            print(f"Error parsing PPTX file {file_path}: {e}")
            return self._create_fallback_report(file_path)
    
    def _extract_structured_content_robust(self, file_path: str) -> Dict[str, Any]:
        """Extract structured content with comprehensive table and chart processing."""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            import pandas as pd
            import io
            
            prs = Presentation(file_path)
            slides_out = []
            
            for si, slide in enumerate(prs.slides, 1):
                slide_dict = {
                    "index": si,
                    "title": slide.shapes.title.text if slide.shapes.title else None,
                    "text_blocks": [],
                    "table_csvs": [],
                    "chart_csvs": [],
                }
                
                # Extract text blocks (skip title shape)
                for shape in slide.shapes:
                    if shape == slide.shapes.title:
                        continue
                    
                    # Extract text from non-structured shapes
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        try:
                            text = ""
                            for paragraph in shape.text_frame.paragraphs:
                                text += paragraph.text + "\n"
                            if text.strip():
                                slide_dict["text_blocks"].append(text.strip())
                        except:
                            pass
                
                # Extract tables and charts as CSV
                for shape in slide.shapes:
                    # Process Table shapes
                    if hasattr(shape, 'has_table') and shape.has_table:
                        try:
                            csv_content = self._extract_table_to_csv(shape.table)
                            if csv_content:
                                slide_dict["table_csvs"].append(csv_content)
                        except Exception as e:
                            print(f"Warning: Failed to extract table from slide {si}: {e}")
                    
                    # Process Chart shapes  
                    elif hasattr(shape, 'has_chart') and shape.has_chart:
                        try:
                            csv_content = self._extract_chart_to_csv(shape.chart)
                            if csv_content:
                                slide_dict["chart_csvs"].append(csv_content)
                        except Exception as e:
                            print(f"Warning: Failed to extract chart from slide {si}: {e}")
                
                slides_out.append(slide_dict)
            
            return {
                "source": str(file_path),
                "num_slides": len(slides_out),
                "slides": slides_out
            }
            
        except Exception as e:
            print(f"Error extracting structured content: {e}")
            return {"slides": []}
    
    def _extract_table_to_csv(self, table) -> str:
        """Convert PowerPoint table to CSV string."""
        try:
            import io
            import csv
            
            # Extract all table data
            data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text_frame.text.replace("\n", " ").strip()
                    row_data.append(cell_text)
                data.append(row_data)
            
            if not data:
                return ""
            
            # Convert to CSV string
            output = io.StringIO()
            writer = csv.writer(output)
            for row in data:
                writer.writerow(row)
            
            csv_content = output.getvalue()
            output.close()
            return csv_content
            
        except Exception as e:
            print(f"Error converting table to CSV: {e}")
            return ""
    
    def _extract_chart_to_csv(self, chart) -> str:
        """Convert PowerPoint chart to CSV string using robust extraction."""
        try:
            import io
            import csv
            
            # Extract categories (x-axis labels)
            categories = []
            try:
                if hasattr(chart, 'plots') and chart.plots:
                    plot = chart.plots[0]
                    if hasattr(plot, 'categories'):
                        for category in plot.categories:
                            if hasattr(category, 'label'):
                                categories.append(str(category.label) if category.label else "")
                            else:
                                categories.append(str(category))
            except:
                pass
            
            # Extract series data
            series_data = []
            if hasattr(chart, 'series'):
                for series in chart.series:
                    series_name = str(series.name) if hasattr(series, 'name') and series.name else "Series"
                    values = []
                    
                    if hasattr(series, 'values'):
                        for value in series.values:
                            try:
                                # Handle different value types
                                if hasattr(value, 'value'):
                                    values.append(value.value)
                                else:
                                    values.append(value)
                            except:
                                values.append("")
                    
                    series_data.append((series_name, values))
            
            if not series_data:
                return ""
            
            # Create CSV with categories as first column and series as subsequent columns
            output = io.StringIO()
            writer = csv.writer(output)
            
            # Create header row
            header = ["Category"] + [series_name for series_name, _ in series_data]
            writer.writerow(header)
            
            # Determine the maximum number of data points
            max_length = max(len(values) for _, values in series_data) if series_data else 0
            if categories:
                max_length = max(max_length, len(categories))
            
            # Write data rows
            for i in range(max_length):
                row = []
                
                # First column: category or row index
                if i < len(categories) and categories[i]:
                    row.append(categories[i])
                else:
                    row.append(f"Row_{i+1}")
                
                # Subsequent columns: series values
                for _, values in series_data:
                    if i < len(values):
                        row.append(values[i])
                    else:
                        row.append("")
                
                writer.writerow(row)
            
            csv_content = output.getvalue()
            output.close()
            return csv_content
            
        except Exception as e:
            print(f"Error converting chart to CSV: {e}")
            return ""
    
    def _process_csv_with_excel_parser(self, csv_content: str, description: str) -> str:
        """Process CSV content with consistent formatting for structured data presentation."""
        try:
            import pandas as pd
            from io import StringIO
            
            # Read CSV content into DataFrame
            df = pd.read_csv(StringIO(csv_content))
            
            if df.empty:
                return f"{description}: No data"
            
            # Format as structured text similar to ExcelParser
            formatted_lines = [f"{description}:"]
            
            # Add formatted table data
            df_str = df.to_string(index=False, na_rep='')
            formatted_lines.append(df_str)
            
            return '\n'.join(formatted_lines)
            
        except Exception as e:
            print(f"Error processing CSV with excel parser: {e}")
            return f"{description}: {csv_content}"
    
    def _create_fallback_report(self, file_path: str) -> Dict[str, Any]:
        """Create minimal report when PPTX parsing fails."""
        filename = Path(file_path).name
        return {
            'metainfo': {
                'sha1_name': filename.rsplit('.', 1)[0],
                'filename': filename,
                'pages_amount': 0,
                'text_blocks_amount': 0,
                'tables_amount': 0,
                'pictures_amount': 0,
                'document_type': 'failed'
            },
            'content': {'pages': []},
            'tables': [],
            'pictures': []
        }


class ExcelParser:
    """Excel file parsing for all sheets and data."""
    
    def parse_excel(self, file_path: str) -> Dict[str, Any]:
        """Parse Excel file for all sheets."""
        try:
            filename = Path(file_path).name
            
            # Read all sheets
            sheets_dict = pd.read_excel(file_path, sheet_name=None, header=None)
            
            pages = []
            total_rows = 0
            
            for sheet_name, df in sheets_dict.items():
                if df.empty:
                    continue
                
                # Convert DataFrame to text representation
                sheet_text_parts = [f"Sheet: {sheet_name}"]
                
                # Convert to string with proper formatting
                df_str = df.to_string(index=False, header=False, na_rep='')
                sheet_text_parts.append(df_str)
                
                combined_sheet_text = '\n'.join(sheet_text_parts)
                
                if combined_sheet_text.strip():
                    pages.append({
                        'page': len(pages) + 1,
                        'text': combined_sheet_text.strip()
                    })
                    total_rows += len(df)
            
            report = {
                'metainfo': {
                    'sha1_name': filename.rsplit('.', 1)[0],
                    'filename': filename,
                    'pages_amount': len(pages),
                    'text_blocks_amount': len(pages),
                    'tables_amount': len(sheets_dict),
                    'pictures_amount': 0,
                    'document_type': 'excel'
                },
                'content': {'pages': pages},
                'tables': [],
                'pictures': []
            }
            
            print(f"Successfully parsed Excel: {filename}")
            print(f"  - {len(pages)} sheets with {total_rows} total rows")
            
            return report
            
        except Exception as e:
            print(f"Error parsing Excel file {file_path}: {e}")
            return self._create_fallback_report(file_path)
    
    def _create_fallback_report(self, file_path: str) -> Dict[str, Any]:
        """Create minimal report when Excel parsing fails."""
        filename = Path(file_path).name
        return {
            'metainfo': {
                'sha1_name': filename.rsplit('.', 1)[0],
                'filename': filename,
                'pages_amount': 0,
                'text_blocks_amount': 0,
                'tables_amount': 0,
                'pictures_amount': 0,
                'document_type': 'failed'
            },
            'content': {'pages': []},
            'tables': [],
            'pictures': []
        }


class UnifiedDocumentParser:
    """Route documents to appropriate parser by file extension."""
    
    def __init__(self):
        self.pdf_parser = PDFParser()
        self.pptx_parser = PPTXParser()
        self.excel_parser = ExcelParser()
    
    def parse_document(self, file_path: str) -> Dict[str, Any]:
        """Parse document based on file extension."""
        try:
            file_ext = Path(file_path).suffix.lower()
            
            if file_ext == '.pdf':
                return self.pdf_parser.parse_pdf(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self.pptx_parser.parse_pptx(file_path)
            elif file_ext in ['.xls', '.xlsx']:
                return self.excel_parser.parse_excel(file_path)
            else:
                print(f"Unsupported file type: {file_ext}")
                return self._create_fallback_report(file_path)
                
        except Exception as e:
            print(f"Error parsing document {file_path}: {e}")
            return self._create_fallback_report(file_path)
    
    def _create_fallback_report(self, file_path: str) -> Dict[str, Any]:
        """Create minimal report when parsing fails."""
        filename = Path(file_path).name
        return {
            'metainfo': {
                'sha1_name': filename.rsplit('.', 1)[0],
                'filename': filename,
                'pages_amount': 0,
                'text_blocks_amount': 0,
                'tables_amount': 0,
                'pictures_amount': 0,
                'document_type': 'failed'
            },
            'content': {'pages': []},
            'tables': [],
            'pictures': []
        }
